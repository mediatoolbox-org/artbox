"""Project initialization module."""

import os

from pathlib import Path
from typing import Any

import yaml

from pptx import Presentation


class FoldedString(str):
    """Custom string class for folded block YAML formatting."""


def folded_string_representer(dumper: yaml.Dumper, data: str) -> yaml.Node:
    """Represent FoldedString using the '>' style."""
    return dumper.represent_scalar("tag:yaml.org,2002:str", data, style=">")


yaml.add_representer(FoldedString, folded_string_representer)


class InitProject:
    """Initialize a new Artbox project from source files."""

    def __init__(self, source_pdf: str, notes_pptx: str, output_path: str):
        """
        Initialize the project scaffolding generator.

        Parameters
        ----------
        source_pdf: str
            Path to the PDF file to use as the visual slides.
        notes_pptx: str
            Path to the PPTX file to use for extracting presenter notes.
        output_path: str
            Path where the generated YAML project configuration will be saved.
        """
        self.source_pdf = source_pdf
        self.notes_pptx = notes_pptx
        self.output_path = output_path

    def _extract_notes(self) -> list[str]:
        """
        Extract presenter notes from each slide in the PPTX.

        Returns
        -------
        list[str]
            A list of extracted notes. Empty string if a slide has no notes.
        """
        prs = Presentation(self.notes_pptx)
        notes = []

        for slide in prs.slides:
            if slide.has_notes_slide:
                text_frame = slide.notes_slide.notes_text_frame
                # text_frame.text will grab all unformatted raw text
                notes.append(text_frame.text.strip())
            else:
                notes.append("")

        return notes

    def generate(self) -> None:
        """
        Generate the project YAML configuration.

        Extracts the notes from the provided PPTX file and constructs a
        compliant dictionary mapping each note string to a page in the
        source PDF.
        """
        if not os.path.exists(self.source_pdf):
            raise FileNotFoundError(f"Source PDF not found: {self.source_pdf}")

        if not os.path.exists(self.notes_pptx):
            raise FileNotFoundError(f"Notes PPTX not found: {self.notes_pptx}")

        # Ensure output directory exists
        os.makedirs(
            os.path.dirname(os.path.abspath(self.output_path)), exist_ok=True
        )

        print(f"Extracting presenter notes from: {self.notes_pptx}")
        notes = self._extract_notes()

        # Calculate relative path from output yaml to the source pdf
        output_dir = Path(self.output_path).parent.resolve()
        pdf_path = Path(self.source_pdf).resolve()

        try:
            rel_pdf_path = os.path.relpath(pdf_path, output_dir)
        except ValueError:
            # Fallback to absolute if on different drives
            rel_pdf_path = str(pdf_path)

        project_name = Path(self.output_path).stem

        # Define the scaffold template matching schema
        scaffold: dict[str, Any] = {
            "name": project_name,
            "output": "/tmp/artbox/",
            "source": {
                "type": "pdf",
                "path": rel_pdf_path,
            },
            "audio": {
                "engine": "edge-tts",
                "instruction": "",
                "defaults": {
                    "language": "en",
                    "gender": "female",
                    "voice-id": "en-US-AriaNeural",
                    "volume": 1.0,
                    "pitch": 1.0,
                    "speed": 1.0,
                },
            },
            "video": {
                "engine": "ffmpeg",
            },
            "slides": {
                "defaults": {
                    "transitions": {
                        "pause-after": 1.0,
                    },
                },
                "items": [],
            },
        }

        # Append each parsed slide
        for i, text in enumerate(notes, start=1):
            audio_text: Any
            if text and len(text.strip()) > 0:
                audio_text = FoldedString(text.strip())
            else:
                audio_text = "Silence."

            slide_config: dict[str, Any] = {
                "slide": i,
                "background": {
                    # PDF pages are 1-indexed for pdf2images
                    "page": i
                },
                "audio": {"text": audio_text},
            }
            scaffold["slides"]["items"].append(slide_config)

        # Output the generated YAML safely
        with open(self.output_path, "w", encoding="utf-8") as file:
            # Use sort_keys=False to maintain logical dictionary ordering
            yaml.dump(
                scaffold,
                file,
                default_flow_style=False,
                sort_keys=False,
                allow_unicode=True,
            )

        print(
            f"Successfully generated project configuration: {self.output_path}"
        )
        print(f"Extracted notes for {len(notes)} slides.")
