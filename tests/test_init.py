"""Tests for the artbox init command."""

import os
from pathlib import Path
import yaml
import pytest

from artbox.init import InitProject
from pptx import Presentation


def test_init_project(tmp_path: Path) -> None:
    """Test the generation of a project YAML from a PPTX and PDF."""
    pdf_path = tmp_path / "presentation.pdf"
    pptx_path = tmp_path / "notes.pptx"
    out_yaml = tmp_path / "my_project.yaml"

    # Create dummy PDF
    pdf_path.touch()

    # Create dummy PPTX with 2 slides (one with notes, one without)
    prs = Presentation()

    # Slide 1: Has notes
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    notes_slide = slide1.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Hello world, this is slide one speaker notes."

    # Slide 2: No notes
    prs.slides.add_slide(prs.slide_layouts[0])

    prs.save(str(pptx_path))

    # Run the init process
    init = InitProject(
        source_pdf=str(pdf_path),
        notes_pptx=str(pptx_path),
        output_path=str(out_yaml),
    )
    init.generate()

    # Assert YAML was created
    assert out_yaml.exists()

    # Verify the contents of the YAML
    with open(out_yaml, "r", encoding="utf-8") as f:
        config = yaml.safe_load(f)

    # 1. Check Root Setup
    assert config["name"] == "my_project"
    assert config["source"]["type"] == "pdf"
    assert config["source"]["path"] == "presentation.pdf"
    assert config["video"]["engine"] == "ffmpeg"
    assert config["audio"]["engine"] == "openai-tts"

    # 2. Check Slide Parsing
    assert len(config["slides"]["items"]) == 2

    # Slide 1 properties
    assert config["slides"]["items"][0]["slide"] == 1
    assert config["slides"]["items"][0]["background"]["page"] == 1
    assert (
        config["slides"]["items"][0]["audio"]["text"]
        == "Hello world, this is slide one speaker notes."
    )

    # Slide 2 properties (fallback to Silence)
    assert config["slides"]["items"][1]["slide"] == 2
    assert config["slides"]["items"][1]["background"]["page"] == 2
    assert config["slides"]["items"][1]["audio"]["text"] == "Silence."
